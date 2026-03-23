"""Document processing Celery task.

Pipeline: Download from R2 → Extract text → Chunk → Embed → Upsert to Qdrant.
"""

import asyncio
from workers.celery_app import celery_app


def _extract_text(content: bytes, filename: str) -> str:
    """Extract text from various document formats.

    Supports: PDF, DOCX, PPTX, TXT, images (via pytesseract).
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        from pypdf import PdfReader
        from io import BytesIO

        reader = PdfReader(BytesIO(content))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)
        return "\n\n".join(pages)

    elif ext == "docx":
        from docx import Document
        from io import BytesIO

        doc = Document(BytesIO(content))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    elif ext == "pptx":
        from pptx import Presentation
        from io import BytesIO

        prs = Presentation(BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts)

    elif ext == "txt":
        return content.decode("utf-8", errors="ignore")

    elif ext in ("png", "jpg", "jpeg", "tiff", "bmp"):
        try:
            import pytesseract
            from PIL import Image
            from io import BytesIO

            img = Image.open(BytesIO(content))
            return pytesseract.image_to_string(img)
        except Exception:
            return ""

    else:
        # Try decoding as UTF-8 text
        try:
            return content.decode("utf-8", errors="ignore")
        except Exception:
            return ""


async def _process_async(file_key: str, user_id: str, institution_id: str, doc_type: str):
    """Async document processing pipeline."""
    from core.storage import storage
    from core.chunker import chunker
    from core.embeddings import embeddings
    from core.vector import qdrant

    # 1. Download from R2
    content = await storage.download(file_key)
    filename = file_key.rsplit("/", 1)[-1]

    # 2. Extract text
    text = _extract_text(content, filename)
    if not text.strip():
        return {"status": "error", "message": "No text extracted"}

    # 3. Chunk text
    chunks = chunker.split_with_metadata(text, base_metadata={
        "user_id": user_id,
        "institution_id": institution_id,
        "doc_type": doc_type,
        "source_file": filename,
        "r2_key": file_key,
    })

    if not chunks:
        return {"status": "error", "message": "No chunks produced"}

    # 4. Embed chunks in batches
    texts_to_embed = [c["text"] for c in chunks]
    vectors = await embeddings.embed_batch(texts_to_embed)

    # 5. Determine target collection
    if doc_type == "courseware":
        collection = f"{institution_id}_courseware"
    elif doc_type == "past_paper":
        collection = f"{institution_id}_{user_id}_exams"
    else:
        collection = f"{institution_id}_{user_id}_notes"

    # 6. Upsert to Qdrant
    points = [
        {"vector": vec, "payload": chunk}
        for vec, chunk in zip(vectors, chunks)
    ]
    count = await qdrant.upsert_batch(collection=collection, points=points)

    # 7. If syllabus, trigger SyllabusMapper
    if doc_type == "syllabus":
        from engines.ace.syllabus_mapper import syllabus_mapper
        await syllabus_mapper.parse(text, institution_id)

    return {"status": "ready", "chunks": count, "collection": collection}


@celery_app.task(name="process_document", bind=True, max_retries=3)
def process_document(self, file_key: str, user_id: str, institution_id: str, doc_type: str = "student_note"):
    """Celery task: process an uploaded document.

    Runs the async pipeline in a new event loop (Celery workers are sync).
    """
    try:
        result = asyncio.run(
            _process_async(file_key, user_id, institution_id, doc_type)
        )
        return result
    except Exception as exc:
        self.retry(exc=exc, countdown=30)
